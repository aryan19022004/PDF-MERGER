import os
import uuid
import base64
import json
from io import BytesIO
import tempfile
import shutil
import zipfile
from flask import Flask, request, send_file, render_template, jsonify, session, redirect, url_for
from pypdf import PdfReader, PdfWriter
import fitz  # PyMuPDF
from PIL import Image
from pdf2docx import Converter as PDF2Docx
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import platform
import subprocess

# Note: comtypes is Windows only
try:
    import comtypes.client
except ImportError:
    pass

app = Flask(__name__)
app.secret_key = 'super-secret-pdf-key'  # Needed for session
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16 MB max limit

# ==========================================
# ADS & ADMIN SYSTEM SETUP (MONGO & CLOUDINARY)
# ==========================================
from pymongo import MongoClient
import cloudinary
import cloudinary.uploader
import bcrypt
from functools import wraps

MONGO_URI = os.environ.get('MONGO_URI', 'mongodb://localhost:27017/pdf_toolkit')
try:
    mongo_client = MongoClient(MONGO_URI, serverSelectionTimeoutMS=5000)
    db = mongo_client.get_database()
    ads_collection = db['ads']
    admin_collection = db['admin']
    
    # Initialize Admin Default Password (@#aryantiwari$%)
    if admin_collection.count_documents({}) == 0:
        default_pw = b"@#aryantiwari$%"
        hashed_pw = bcrypt.hashpw(default_pw, bcrypt.gensalt())
        admin_collection.insert_one({"role": "admin", "password": hashed_pw})
        print("Default Admin Credentials Generated.")
except Exception as e:
    print(f"MongoDB Configuration Warning: {e}")

cloudinary.config(
    cloud_name = os.environ.get('CLOUDINARY_CLOUD_NAME', 'YOUR_CLOUD_NAME'),
    api_key = os.environ.get('CLOUDINARY_API_KEY', 'YOUR_API_KEY'),
    api_secret = os.environ.get('CLOUDINARY_API_SECRET', 'YOUR_API_SECRET'),
    secure = True
)

def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not session.get('admin_logged_in'):
            return redirect(url_for('admin_login'))
        return f(*args, **kwargs)
    return decorated_function
# ==========================================

# Dictionary to hold active editing sessions temporarily in memory
EDIT_SESSIONS = {}

def parse_page_range(range_str, total_pages):
    """
    Parses a string like '1-3, 5, 7-9' into a list of 0-based page indices.
    If range_str is empty or None, returns all pages.
    """
    if not range_str or not str(range_str).strip():
        return list(range(total_pages))
    
    pages = set()
    parts = str(range_str).replace(' ', '').split(',')
    
    try:
        for part in parts:
            if not part:
                continue
            if '-' in part:
                start_str, end_str = part.split('-', 1)
                start = int(start_str) - 1
                end = int(end_str) - 1
                
                start = max(0, start)
                end = min(total_pages - 1, end)
                
                if start <= end:
                    pages.update(range(start, end + 1))
            else:
                page = int(part) - 1
                if 0 <= page < total_pages:
                    pages.add(page)
    except ValueError as e:
        raise ValueError("Invalid range format. Use numbers, commas, and dashes.")
        
    return sorted(list(pages))

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/merge')
def merge_page():
    return render_template('merge.html')

@app.route('/split', methods=['GET', 'POST'])
def split_page():
    if request.method == 'POST':
        file = request.files.get('pdf_file')
        if not file or file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
            
        try:
            reader = PdfReader(file)
            total_pages = len(reader.pages)
            pages_str = request.form.get('pages', '').strip()
            
            if pages_str:
                # Extract specific pages into one PDF
                selected_pages = parse_page_range(pages_str, total_pages)
                writer = PdfWriter()
                for p in selected_pages:
                    writer.add_page(reader.pages[p])
                    
                output_pdf = BytesIO()
                writer.write(output_pdf)
                output_pdf.seek(0)
                
                return send_file(
                    output_pdf,
                    mimetype='application/pdf',
                    as_attachment=True,
                    download_name=f"split_{file.filename}"
                )
            else:
                # Extract all individual pages into a ZIP
                outputs = []
                for p in range(total_pages):
                    writer = PdfWriter()
                    writer.add_page(reader.pages[p])
                    out_io = BytesIO()
                    writer.write(out_io)
                    outputs.append((f"page_{p+1}.pdf", out_io.getvalue()))
                    
                return serve_converted_files(outputs, f"split_{os.path.splitext(file.filename)[0]}.zip")
                
        except Exception as e:
            return jsonify({'error': str(e)}), 400
            
    return render_template('split.html')

@app.route('/edit')
def edit_page():
    return render_template('edit.html')

@app.route('/pdf-to-excel')
def pdf_to_excel():
    return render_template('pdf_to_excel.html')

@app.route('/excel-to-pdf')
def excel_to_pdf():
    return render_template('excel_to_pdf.html')

@app.route('/pdf-to-word')
def pdf_to_word():
    return render_template('pdf_to_word.html')

@app.route('/word-to-pdf')
def word_to_pdf():
    return render_template('word_to_pdf.html')

@app.route('/pdf-to-ppt')
def pdf_to_ppt():
    return render_template('pdf_to_ppt.html')

@app.route('/ppt-to-pdf')
def ppt_to_pdf():
    return render_template('ppt_to_pdf.html')

@app.route('/jpg-to-pdf')
def jpg_to_pdf():
    return render_template('jpg_to_pdf.html')

@app.route('/pdf-to-jpg')
def pdf_to_jpg():
    return render_template('pdf_to_jpg.html')

@app.route('/protect-pdf')
def protect_pdf():
    return render_template('protect_pdf.html')

@app.route('/unprotect-pdf')
def unprotect_pdf():
    return render_template('unprotect_pdf.html')

@app.route('/pdf-to-ocr')
def pdf_to_ocr():
    return render_template('pdf_to_ocr.html')

@app.route('/ocr-to-pdf')
def ocr_to_pdf():
    return render_template('ocr_to_pdf.html')

@app.route('/watermark')
def watermark_page():
    return render_template('watermark.html')

@app.route('/remove-pages')
def remove_pages_page():
    return render_template('remove_pages.html')

@app.route('/add-pages')
def add_pages_page():
    return render_template('add_pages.html')

@app.route('/rearrange')
def rearrange_page():
    return render_template('rearrange.html')

# ==========================================
# ADMIN & ADS ROUTES
# ==========================================
from datetime import datetime, timedelta
from bson.objectid import ObjectId

@app.route('/admin')
def admin_route():
    if session.get('admin_logged_in'):
        return redirect(url_for('admin_dashboard'))
    return render_template('admin_login.html')

@app.route('/api/admin/login', methods=['POST'])
def api_admin_login():
    try:
        data = request.json
        password = data.get('password', '').encode('utf-8')
        
        admin_doc = admin_collection.find_one({"role": "admin"})
        if not admin_doc:
            return jsonify({"error": "Admin account not found"}), 500
            
        if bcrypt.checkpw(password, admin_doc['password']):
            session['admin_logged_in'] = True
            return jsonify({"success": True})
        else:
            return jsonify({"error": "Invalid password"}), 401
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/admin/dashboard')
@admin_required
def admin_dashboard():
    return render_template('admin_dashboard.html')

@app.route('/api/admin/change-password', methods=['POST'])
@admin_required
def api_change_password():
    try:
        data = request.json
        new_password = data.get('new_password', '').encode('utf-8')
        if len(new_password) < 6:
            return jsonify({"error": "Password too short"}), 400
            
        hashed_pw = bcrypt.hashpw(new_password, bcrypt.gensalt())
        admin_collection.update_one({"role": "admin"}, {"$set": {"password": hashed_pw}})
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/ads/create', methods=['POST'])
@admin_required
def api_create_ad():
    try:
        title = request.form.get('title', 'Untitled Ad')
        redirect_link = request.form.get('redirectLink', '#')
        duration_days = int(request.form.get('durationDays', 7))
        skip_seconds = int(request.form.get('skipAfterSeconds', 5))
        
        file = request.files.get('image')
        if not file:
            return jsonify({"error": "Image is required"}), 400
            
        # Upload to Cloudinary
        upload_result = cloudinary.uploader.upload(file)
        image_url = upload_result.get('secure_url')
        
        now = datetime.utcnow()
        expires_at = now + timedelta(days=duration_days)
        
        ad_data = {
            "title": title,
            "imageUrl": image_url,
            "redirectLink": redirect_link,
            "durationDays": duration_days,
            "skipAfterSeconds": skip_seconds,
            "createdAt": now,
            "expiresAt": expires_at,
            "viewsCount": 0,
            "totalWatchTime": 0,
            "clicks": 0
        }
        
        ads_collection.insert_one(ad_data)
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/ads/list', methods=['GET'])
@admin_required
def api_list_ads():
    try:
        ads = list(ads_collection.find().sort("createdAt", -1))
        for ad in ads:
            ad['_id'] = str(ad['_id'])
        return jsonify({"ads": ads})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/ads/delete/<ad_id>', methods=['DELETE'])
@admin_required
def api_delete_ad(ad_id):
    try:
        ads_collection.delete_one({"_id": ObjectId(ad_id)})
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# Public APIs for Ad Display
@app.route('/api/ads/random', methods=['GET'])
def api_random_ad():
    try:
        # Only fetch ads that haven't expired
        now = datetime.utcnow()
        pipeline = [
            {"$match": {"expiresAt": {"$gt": now}}},
            {"$sample": {"size": 1}}
        ]
        random_ads = list(ads_collection.aggregate(pipeline))
        if not random_ads:
            return jsonify({"ad": None})
            
        ad = random_ads[0]
        ad['_id'] = str(ad['_id'])
        return jsonify({"ad": {"_id": ad['_id'], "imageUrl": ad['imageUrl'], "redirectLink": ad['redirectLink'], "skipAfterSeconds": ad['skipAfterSeconds']}})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/ads/track', methods=['POST'])
def api_track_ad():
    try:
        data = request.json
        ad_id = data.get('ad_id')
        action = data.get('action') # 'view' or 'click'
        watch_time = data.get('watch_time', 0)
        
        if not ad_id:
            return jsonify({"error": "ad_id required"}), 400
            
        update_query = {}
        if action == 'view':
            update_query = {"$inc": {"viewsCount": 1, "totalWatchTime": watch_time}}
        elif action == 'click':
            update_query = {"$inc": {"clicks": 1, "totalWatchTime": watch_time}}
            
        if update_query:
            ads_collection.update_one({"_id": ObjectId(ad_id)}, update_query)
            
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ==========================================
# TOOL ENDPOINTS 
# ==========================================

@app.route('/merge', methods=['POST'])
def merge_pdfs():
    try:
        indices = []
        for key in request.files.keys():
            if key.startswith('pdf_'):
                try:
                    idx = int(key.split('_')[1])
                    indices.append(idx)
                except ValueError:
                    pass
                    
        indices.sort()
        
        if len(indices) < 2:
            return jsonify({'error': 'At least two PDF files are required'}), 400
            
        writer = PdfWriter()
        
        for idx in indices:
            file = request.files[f'pdf_{idx}']
            if file.filename == '':
                return jsonify({'error': 'All blocks must have a PDF file selected'}), 400
                
            if not file.filename.lower().endswith('.pdf'):
                return jsonify({'error': f'File {file.filename} is not a valid PDF'}), 400
                
            reader = PdfReader(file)
            total_pages = len(reader.pages)
            
            range_str = request.form.get(f'range_{idx}', '')
            
            try:
                pages = parse_page_range(range_str, total_pages)
            except ValueError as e:
                return jsonify({'error': f"Error in document {idx+1}: {str(e)}"}), 400
                
            for p in pages:
                if 0 <= p < total_pages:
                    writer.add_page(reader.pages[p])

        output_pdf = BytesIO()
        writer.write(output_pdf)
        output_pdf.seek(0)
        
        return send_file(
            output_pdf,
            mimetype='application/pdf',
            as_attachment=True,
            download_name='merged_output.pdf'
        )
        
    except Exception as e:
        return jsonify({'error': f"An error occurred: {str(e)}"}), 500

@app.route('/api/protect', methods=['POST'])
def api_protect():
    file = request.files.get('pdf_file')
    password = request.form.get('password')
    if not file or not password:
        return jsonify({'error': 'File and password are required'}), 400
    try:
        pdf_data = file.read()
        doc = fitz.open(stream=pdf_data, filetype="pdf")
        
        # Save to buffer with encryption
        out_buf = BytesIO()
        doc.save(out_buf, encryption=fitz.PDF_ENCRYPT_AES_256, user_pw=password, owner_pw=password)
        out_buf.seek(0)
        
        return send_file(
            out_buf,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=f"Protected_{file.filename}"
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/unprotect', methods=['POST'])
def api_unprotect():
    file = request.files.get('pdf_file')
    password = request.form.get('password', '')
    if not file:
        return jsonify({'error': 'File is required'}), 400
    try:
        pdf_data = file.read()
        doc = fitz.open(stream=pdf_data, filetype="pdf")
        
        if doc.needs_pass:
            if not doc.authenticate(password):
                return jsonify({'error': 'Incorrect password'}), 403
        
        # Save decrypted
        out_buf = BytesIO()
        doc.save(out_buf)
        out_buf.seek(0)
        
        return send_file(
            out_buf,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=f"Unlocked_{file.filename}"
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/watermark', methods=['POST'])
def api_watermark():
    file = request.files.get('pdf_file')
    text = request.form.get('text', 'Watermark')
    position = request.form.get('position', 'center') # top, center, bottom
    opacity = float(request.form.get('opacity', 0.5))
    size = int(request.form.get('size', 48))
    
    if not file:
        return jsonify({'error': 'File is required'}), 400
        
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        for page in doc:
            rect = page.rect
            x = rect.width / 2.0
            
            if position == 'top':
                y = rect.height * 0.15
            elif position == 'bottom':
                y = rect.height * 0.85
            else:
                y = rect.height / 2.0
                
            text_length = fitz.get_text_length(text, fontsize=size)
            x_pos = (rect.width - text_length) / 2.0
            
            # Using basic insert_text
            page.insert_text(fitz.Point(x_pos, y), text, fontsize=size, fill_opacity=opacity, color=(0.5, 0.5, 0.5))
            
        out_buf = BytesIO()
        doc.save(out_buf)
        out_buf.seek(0)
        return send_file(out_buf, mimetype='application/pdf', as_attachment=True, download_name=f"Watermarked_{file.filename}")
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/remove-pages', methods=['POST'])
def api_remove_pages():
    file = request.files.get('pdf_file')
    pages_str = request.form.get('pages', '')
    if not file or not pages_str:
        return jsonify({'error': 'File and pages to remove are required'}), 400
        
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        total = len(doc)
        
        # Determine pages to remove
        to_remove = set()
        for part in pages_str.replace(' ', '').split(','):
            if '-' in part:
                s, e = part.split('-')
                to_remove.update(range(int(s)-1, int(e)))
            else:
                to_remove.add(int(part)-1)
                
        to_keep = [i for i in range(total) if i not in to_remove and 0 <= i < total]
        if not to_keep:
            return jsonify({'error': 'Cannot remove all pages'}), 400
            
        doc.select(to_keep)
        out_buf = BytesIO()
        doc.save(out_buf)
        out_buf.seek(0)
        return send_file(out_buf, mimetype='application/pdf', as_attachment=True, download_name=f"Reduced_{file.filename}")
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/add-pages', methods=['POST'])
def api_add_pages():
    main_file = request.files.get('main_pdf')
    addon_file = request.files.get('addon_file')
    position = request.form.get('position', 'end') # start, end, custom
    custom_page = int(request.form.get('custom_page', 1)) - 1
    
    if not main_file or not addon_file:
        return jsonify({'error': 'Main PDF and Addon file are required'}), 400
        
    try:
        main_doc = fitz.open(stream=main_file.read(), filetype="pdf")
        addon_data = addon_file.read()
        
        addon_doc = fitz.open(stream=addon_data, filetype="pdf" if addon_file.filename.lower().endswith('.pdf') else None)
        
        # If it's an image, fitz can open it, but we convert it to PDF
        if not addon_file.filename.lower().endswith('.pdf'):
            pdfbytes = addon_doc.convert_to_pdf()
            addon_doc = fitz.open("pdf", pdfbytes)
            
        insert_idx = len(main_doc)
        if position == 'start':
            insert_idx = 0
        elif position == 'custom':
            insert_idx = max(0, min(custom_page, len(main_doc)))
            
        main_doc.insert_pdf(addon_doc, start_at=insert_idx)
        
        out_buf = BytesIO()
        main_doc.save(out_buf)
        out_buf.seek(0)
        return send_file(out_buf, mimetype='application/pdf', as_attachment=True, download_name=f"Expanded_{main_file.filename}")
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/rearrange', methods=['POST'])
def api_rearrange():
    file = request.files.get('pdf_file')
    order_str = request.form.get('order') # e.g. "2,0,1" -> representing 0-based indices
    if not file or not order_str:
        return jsonify({'error': 'File and ordering are required'}), 400
        
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        order = [int(x.strip()) for x in order_str.split(',')]
        
        # Validate order length and bounds
        if len(order) != len(doc):
            pass # We could permit skipping pages, but usually rearrange means all pages
            
        valid_order = [p for p in order if 0 <= p < len(doc)]
        
        doc.select(valid_order)
        out_buf = BytesIO()
        doc.save(out_buf)
        out_buf.seek(0)
        return send_file(out_buf, mimetype='application/pdf', as_attachment=True, download_name=f"Rearranged_{file.filename}")
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ==========================================
# ==========================================

def serve_converted_files(outputs, zip_name):
    if not outputs:
        return jsonify({'error': 'No completed conversions to return'}), 400
    if len(outputs) == 1:
        out_io = BytesIO(outputs[0][1])
        return send_file(out_io, as_attachment=True, download_name=outputs[0][0])
    zip_buffer = BytesIO()
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for fname, b_content in outputs:
            zipf.writestr(fname, b_content)
    zip_buffer.seek(0)
    return send_file(zip_buffer, mimetype="application/zip", as_attachment=True, download_name=zip_name)

@app.route('/api/convert/pdf-to-word', methods=['POST'])
def api_pdf_to_word():
    files = request.files.getlist('files')
    outputs = []
    with tempfile.TemporaryDirectory() as tmpdir:
        for f in files:
            if not f.filename: continue
            pdf_path = os.path.join(tmpdir, f.filename)
            docx_name = os.path.splitext(f.filename)[0] + ".docx"
            docx_path = os.path.join(tmpdir, docx_name)
            f.save(pdf_path)
            try:
                cv = PDF2Docx(pdf_path)
                cv.convert(docx_path)
                cv.close()
                with open(docx_path, 'rb') as df:
                    outputs.append((docx_name, df.read()))
            except Exception as e:
                print(e)
    return serve_converted_files(outputs, "word_docs.zip")

@app.route('/api/convert/word-to-pdf', methods=['POST'])
def api_word_to_pdf():
    files = request.files.getlist('files')
    outputs = []
    with tempfile.TemporaryDirectory() as tmpdir:
        for f in files:
            if not f.filename: continue
            docx_path = os.path.join(tmpdir, f.filename)
            pdf_name = os.path.splitext(f.filename)[0] + ".pdf"
            pdf_path = os.path.join(tmpdir, pdf_name)
            f.save(docx_path)
            try:
                # Need to use abspath for docx2pdf
                abs_in = os.path.abspath(docx_path)
                abs_out = os.path.abspath(pdf_path)
                
                if platform.system() == 'Windows':
                    from docx2pdf import convert as docx2pdf_convert
                    docx2pdf_convert(abs_in, abs_out)
                else:
                    subprocess.run(['libreoffice', '--headless', '--nologo', '--convert-to', 'pdf', '--outdir', tmpdir, abs_in], check=True)
                
                with open(abs_out, 'rb') as pf:
                    outputs.append((pdf_name, pf.read()))
            except Exception as e:
                print("Error converting Word to PDF:", e)
    return serve_converted_files(outputs, "converted_pdfs.zip")

@app.route('/api/convert/pdf-to-excel', methods=['POST'])
def api_pdf_to_excel():
    import pdfplumber
    files = request.files.getlist('files')
    outputs = []
    for f in files:
        if not f.filename: continue
        try:
            excel_name = os.path.splitext(f.filename)[0] + ".xlsx"
            out_io = BytesIO()
            with pdfplumber.open(f) as pdf:
                writer = pd.ExcelWriter(out_io, engine='openpyxl')
                for i, page in enumerate(pdf.pages):
                    table = page.extract_table()
                    if table:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        df.to_excel(writer, sheet_name=f"Page_{i+1}", index=False)
                writer.close()
            outputs.append((excel_name, out_io.getvalue()))
        except Exception as e:
            print(e)
    return serve_converted_files(outputs, "excel_sheets.zip")

@app.route('/api/convert/excel-to-pdf', methods=['POST'])
def api_excel_to_pdf():
    # Simplistic conversion: Excel -> HTML -> text for PDF
    files = request.files.getlist('files')
    outputs = []
    for f in files:
        if not f.filename: continue
        try:
            pdf_name = os.path.splitext(f.filename)[0] + ".pdf"
            df = pd.read_excel(f)
            doc = fitz.open()
            page = doc.new_page()
            text = df.to_string()
            page.insert_text(fitz.Point(50, 50), text, fontsize=8)
            out_io = BytesIO()
            doc.save(out_io)
            outputs.append((pdf_name, out_io.getvalue()))
        except Exception as e:
            print(e)
    return serve_converted_files(outputs, "converted_pdfs.zip")

@app.route('/api/convert/pdf-to-ppt', methods=['POST'])
def api_pdf_to_ppt():
    files = request.files.getlist('files')
    outputs = []
    for f in files:
        if not f.filename: continue
        try:
            ppt_name = os.path.splitext(f.filename)[0] + ".pptx"
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            pdf_bytes = f.read()
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            for page in doc:
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_stream = BytesIO(pix.tobytes("png"))
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(img_stream, 0, 0, width=Inches(10))
            out_ppt = BytesIO()
            prs.save(out_ppt)
            outputs.append((ppt_name, out_ppt.getvalue()))
        except Exception as e:
            print(e)
    return serve_converted_files(outputs, "presentations.zip")

@app.route('/api/convert/ppt-to-pdf', methods=['POST'])
def api_ppt_to_pdf():
    files = request.files.getlist('files')
    outputs = []
    with tempfile.TemporaryDirectory() as tmpdir:
        for f in files:
            if not f.filename: continue
            ppt_name = f.filename
            ppt_path = os.path.join(tmpdir, ppt_name)
            pdf_name = os.path.splitext(ppt_name)[0] + ".pdf"
            pdf_path = os.path.join(tmpdir, pdf_name)
            f.save(ppt_path)
            try:
                abs_in = os.path.abspath(ppt_path)
                abs_out = os.path.abspath(pdf_path)
                
                if platform.system() == 'Windows':
                    import comtypes.client
                    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                    # Headless presentation save
                    deck = powerpoint.Presentations.Open(abs_in, WithWindow=False)
                    deck.SaveAs(abs_out, 32) # 32 is ppSaveAsPDF
                    deck.Close()
                    powerpoint.Quit()
                else:
                    subprocess.run(['libreoffice', '--headless', '--nologo', '--convert-to', 'pdf', '--outdir', tmpdir, abs_in], check=True)
                
                with open(abs_out, 'rb') as pf:
                    outputs.append((pdf_name, pf.read()))
            except Exception as e:
                print("Error converting PPT to PDF:", e)
    return serve_converted_files(outputs, "converted_pdfs.zip")

@app.route('/api/convert/jpg-to-pdf', methods=['POST'])
def api_jpg_to_pdf():
    files = request.files.getlist('files')
    images = []
    for f in files:
        if f.filename.lower().endswith(('.png', '.jpg', '.jpeg', '.tiff')):
            try:
                img = Image.open(f).convert('RGB')
                images.append(img)
            except Exception as e:
                print(e)
    if not images:
        return jsonify({'error': 'No valid images found'}), 400
    out_pdf = BytesIO()
    images[0].save(out_pdf, format='PDF', save_all=True, append_images=images[1:])
    out_pdf.seek(0)
    return send_file(out_pdf, as_attachment=True, download_name='converted_images.pdf')

@app.route('/api/convert/pdf-to-jpg', methods=['POST'])
def api_pdf_to_jpg():
    files = request.files.getlist('files')
    outputs = []
    for f in files:
        if not f.filename: continue
        try:
            base_name = os.path.splitext(f.filename)[0]
            pdf_bytes = f.read()
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                jpg_name = f"{base_name}_page_{i+1}.jpg"
                outputs.append((jpg_name, pix.tobytes("jpeg")))
        except Exception as e:
            print(e)
    return serve_converted_files(outputs, "extracted_images.zip")

# ==========================================
# PDF EDITOR ROUTES (PHASE 2)
# ==========================================

@app.route('/api/upload_edit', methods=['POST'])
def upload_edit():
    """Uploads a PDF and creates a memory session."""
    if 'edit_pdf' not in request.files:
        return jsonify({'error': 'PDF file is required'}), 400
        
    file = request.files['edit_pdf']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
        
    session_id = str(uuid.uuid4())
    pdf_bytes = file.read()
    
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        total_pages = len(doc)
    except Exception as e:
        return jsonify({'error': 'Invalid PDF file'}), 400
        
    EDIT_SESSIONS[session_id] = {
        'bytes': pdf_bytes,
        'filename': file.filename,
        'doc': doc
    }
    
    return jsonify({
        'session_id': session_id,
        'total_pages': total_pages
    })

@app.route('/api/get_pdf/<session_id>')
def get_pdf(session_id):
    """Returns the raw PDF bytes for PDF.js to render."""
    if session_id not in EDIT_SESSIONS:
        return jsonify({'error': 'Session expired or invalid'}), 404
        
    pdf_bytes = EDIT_SESSIONS[session_id]['bytes']
    return send_file(
        BytesIO(pdf_bytes),
        mimetype='application/pdf'
    )

@app.route('/api/page_data/<session_id>/<int:page_num>')
def get_page_data(session_id, page_num):
    """Returns the background image and text blocks for a specific page."""
    if session_id not in EDIT_SESSIONS:
        return jsonify({'error': 'Session expired or invalid'}), 404
        
    doc = EDIT_SESSIONS[session_id]['doc']
    
    if page_num < 0 or page_num >= len(doc):
        return jsonify({'error': 'Invalid page number'}), 400
        
    page = doc[page_num]
    
    # 1. Render page to image for the canvas background
    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2)) # 2x zoom for better resolution
    img_data = pix.tobytes("png")
    img_base64 = base64.b64encode(img_data).decode('utf-8')
    
    # 2. Extract text blocks (words/lines) and their coordinates
    # words = page.get_text("words")  # format: [x0, y0, x1, y1, "word", block_no, line_no, word_no]
    dict_data = page.get_text("dict")
    text_blocks = []
    
    # Scale coordinates to match the HTML canvas space based on image vs original rect
    orig_rect = page.rect
    
    for block in dict_data.get("blocks", []):
        if block.get("type") == 0:  # Text block
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    if span.get("text", "").strip():
                        # Extract BBox: [x0, y0, x1, y1]
                        bbox = span["bbox"]
                        text_blocks.append({
                            'id': str(uuid.uuid4()),
                            'text': span["text"],
                            'bbox': bbox,
                            'font': span["font"],
                            'size': span["size"],
                            'color': span["color"],  # int representation
                            'flags': span["flags"]   # bold, italic, etc.
                        })
                        
    return jsonify({
        'image': f"data:image/png;base64,{img_base64}",
        'width': orig_rect.width,
        'height': orig_rect.height,
        'blocks': text_blocks
    })

@app.route('/api/reorder_pages', methods=['POST'])
def reorder_pages():
    data = request.json
    session_id = data.get('session_id')
    new_order = data.get('new_order', [])  # list of ints or "BLANK"
    
    if session_id not in EDIT_SESSIONS:
        return jsonify({'error': 'Session expired'}), 404
        
    doc = EDIT_SESSIONS[session_id]['doc']
    new_doc = fitz.open()
    
    for item in new_order:
        if item == 'BLANK':
            w, h = 595, 842
            if len(doc) > 0:
                rect = doc[0].rect
                w, h = rect.width, rect.height
            new_doc.new_page(width=w, height=h)
        else:
            idx = int(item)
            if 0 <= idx < len(doc):
                new_doc.insert_pdf(doc, from_page=idx, to_page=idx)
                
    EDIT_SESSIONS[session_id]['doc'] = new_doc
    
    out_pdf = BytesIO()
    new_doc.save(out_pdf, garbage=4, deflate=True)
    EDIT_SESSIONS[session_id]['bytes'] = out_pdf.getvalue()
    
    return jsonify({'success': True, 'total_pages': len(new_doc)})

@app.route('/api/delete_page', methods=['POST'])
def delete_page():
    data = request.json
    session_id = data.get('session_id')
    page_num = data.get('page_num')
    
    if session_id not in EDIT_SESSIONS:
        return jsonify({'error': 'Session expired'}), 404
        
    doc = EDIT_SESSIONS[session_id]['doc']
    
    if page_num < 0 or page_num >= len(doc):
        return jsonify({'error': 'Invalid page number'}), 400
        
    doc.delete_page(page_num)
    
    out_pdf = BytesIO()
    doc.save(out_pdf, garbage=4, deflate=True)
    EDIT_SESSIONS[session_id]['bytes'] = out_pdf.getvalue()
    
    return jsonify({'success': True, 'total_pages': len(doc)})

@app.route('/page-number')
def page_number():
    return render_template('page_number.html')

@app.route('/api/add_page_numbers', methods=['POST'])
def add_page_numbers():
    if 'pdf_file' not in request.files:
        return jsonify({'error': 'PDF file is required'}), 400
        
    file = request.files['pdf_file']
    range_str = request.form.get('range', '').strip()
    prefix = request.form.get('prefix', '')
    try:
        start_num = int(request.form.get('start_num', 1))
    except:
        start_num = 1
    position = request.form.get('position', 'bottom-center')
    try:
        size = int(request.form.get('size', 12))
    except:
        size = 12
    color_hex = request.form.get('color', '#000000').lstrip('#')
    
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
    except Exception as e:
        return jsonify({'error': 'Invalid PDF file'}), 400
        
    total_pages = len(doc)
    
    if not range_str:
        pages = list(range(total_pages))
    else:
        try:
            pages = parse_page_range(range_str, total_pages)
        except ValueError as e:
            return jsonify({'error': str(e)}), 400
            
    try:
        r = int(color_hex[0:2], 16) / 255.0
        g = int(color_hex[2:4], 16) / 255.0
        b = int(color_hex[4:6], 16) / 255.0
        color_tuple = (r, g, b)
    except:
        color_tuple = (0, 0, 0)

    for p in pages:
        if 0 <= p < total_pages:
            page = doc[p]
            rect = page.rect
            text_str = f"{prefix}{start_num}"
            start_num += 1
            
            # Rough width approximation: character count * fontsize * 0.5
            apx_width = len(text_str) * size * 0.5
            
            x = rect.width / 2.0 - apx_width / 2.0
            y = rect.height - 30
            
            if position == 'bottom-right':
                x = rect.width - 30 - apx_width
            elif position == 'top-right':
                x = rect.width - 30 - apx_width
                y = 30 + size
                
            point = fitz.Point(x, y)
            page.insert_text(point, text_str, fontname="helv", fontsize=size, color=color_tuple)
            
    out_pdf = BytesIO()
    doc.save(out_pdf, garbage=4, deflate=True)
    out_pdf.seek(0)
    
    return send_file(out_pdf, as_attachment=True, download_name='numbered_document.pdf', mimetype='application/pdf')

@app.route('/api/save_edit', methods=['POST'])
def save_edit():
    """Applies changes from the frontend to the PDF document."""
    data = request.json
    session_id = data.get('session_id')
    page_num = data.get('page_num')
    edits = data.get('edits', []) # List of edited blocks
    deleted_blocks = data.get('deleted', []) # Original blocks that were modified
    
    if session_id not in EDIT_SESSIONS:
        return jsonify({'error': 'Session expired or invalid'}), 404
        
    doc = EDIT_SESSIONS[session_id]['doc']
    page = doc[page_num]
    
    # 1. Erase original text by drawing a white rectangle over the old bounds
    # This is not perfect text replacement but standard for overlaid PDF edits
    for del_bbox in deleted_blocks:
        rect = fitz.Rect(del_bbox)
        # Redact/erase
        page.add_redact_annot(rect, fill=(1, 1, 1)) 
        
    page.apply_redactions()
    
    # 2. Insert new text at specified coordinates
    for edit in edits:
        text = edit['text']
        bbox = edit['new_bbox'] # [x0, y0, x1, y1]
        fontname = edit.get('fontFamilly', 'helv') 
        fontsize = edit.get('fontSize', 12)
        
        # Parse hex color back to rgb floats
        color_hex = edit.get('color', '#000000').lstrip('#')
        try:
            r = int(color_hex[0:2], 16) / 255.0
            g = int(color_hex[2:4], 16) / 255.0
            b = int(color_hex[4:6], 16) / 255.0
            color_tuple = (r, g, b)
        except:
            color_tuple = (0, 0, 0)
            
        try:
            # We attempt to insert the text with native spacing relying on HTML logic
            # This handles Bold, Italic, and normal spans exactly as they looked
            css = f"font-family: {fontname}; font-size: {fontsize}px; color: #{color_hex};"
            html = edit.get('html', text)
            if not html.startswith('<'):
                html = f"<div style='{css}'>{html}</div>"
            else:
                html = f"<div style='{css}'>{html}</div>"
                
            # rect [x0, y0, x1, y1] given by tracking code
            # We expand the right/bottom slightly to prevent word-wrapping constraints
            y_offset = float(fontsize) * 0.2
            rect = fitz.Rect(bbox[0], bbox[1], bbox[2] + 200, bbox[3] + y_offset + 50)
            page.insert_htmlbox(rect, html, css=css)
        except Exception as e:
            print(f"Error inserting: {e}")
            
    # Save the current state in memory
    EDIT_SESSIONS[session_id]['doc'] = doc
    return jsonify({'success': True})

@app.route('/api/download_edit/<session_id>')
def download_edit(session_id):
    """Downloads the edited PDF."""
    if session_id not in EDIT_SESSIONS:
        return "Not found", 404
        
    doc = EDIT_SESSIONS[session_id]['doc']
    filename = EDIT_SESSIONS[session_id]['filename']
    
    out_pdf = BytesIO()
    doc.save(out_pdf, garbage=4, deflate=True)
    out_pdf.seek(0)
    
    # We can delete the session if we assume download = finish
    # del EDIT_SESSIONS[session_id]
    
    return send_file(
        out_pdf,
        mimetype='application/pdf',
        as_attachment=True,
        download_name=f"edited_{filename}"
    )

if __name__ == '__main__':
    app.run(debug=True, port=5000)
